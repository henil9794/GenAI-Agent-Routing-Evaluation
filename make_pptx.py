from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

RESULTS = "/Users/parinshah/Documents/NEU/Sem IV/GenAI/Project/GenAI-Agent-Routing-Evaluation/data/results"
OUT = "/Users/parinshah/Documents/NEU/Sem IV/GenAI/Project/GenAI-Agent-Routing-Evaluation/data/results/team8_presentation.pptx"

WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG   = RGBColor(0xF7, 0xF9, 0xFC)
DARK_TEXT  = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT     = RGBColor(0x1A, 0x73, 0xE8)
ACCENT2    = RGBColor(0x0F, 0x9D, 0x58)
ACCENT3    = RGBColor(0xF4, 0xB4, 0x00)
LIGHT_LINE = RGBColor(0xDD, 0xE3, 0xEA)
SUBTITLE   = RGBColor(0x5F, 0x6B, 0x7C)
CARD_BG    = RGBColor(0xEB, 0xF3, 0xFD)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

def add_slide():
    s = prs.slides.add_slide(blank)
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = LIGHT_BG
    return s

def rect(slide, l, t, w, h, fill=WHITE, line=None):
    shp = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line:
        shp.line.color.rgb = line; shp.line.width = Pt(1)
    else:
        shp.line.fill.background()
    return shp

def txt(slide, text, l, t, w, h, size=18, bold=False, color=DARK_TEXT,
        align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    box.word_wrap = True
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.size = Pt(size); run.font.bold = bold
    run.font.italic = italic; run.font.color.rgb = color
    return box

def add_img(slide, path, l, t, w, h=None):
    if h:
        slide.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h))
    else:
        slide.shapes.add_picture(path, Inches(l), Inches(t), Inches(w))

def accent_bar(slide, color=ACCENT):
    shp = slide.shapes.add_shape(1, Inches(0.5), Inches(0.38), Inches(0.06), Inches(0.55))
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    shp.line.fill.background()

def slide_num(slide, n, total=6):
    txt(slide, f"{n} / {total}", 12.3, 7.1, 0.8, 0.3, size=9, color=SUBTITLE, align=PP_ALIGN.RIGHT)

def color_bar(slide, x, y, w, h, color):
    b = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    b.fill.solid(); b.fill.fore_color.rgb = color; b.line.fill.background()

# ── SLIDE 1: Title ──────────────────────────────────────────────────────────
s1 = add_slide()
color_bar(s1, 0, 0, 13.33, 0.12, ACCENT)
rect(s1, 1.5, 1.4, 10.33, 4.7, fill=WHITE, line=LIGHT_LINE)
color_bar(s1, 1.5, 1.4, 0.12, 4.7, ACCENT)
txt(s1, "Agentic RAG:", 2.0, 2.0, 9.5, 0.9, size=40, bold=True)
txt(s1, "Tool Routing Under Prompt Ambiguity", 2.0, 2.85, 9.5, 0.7, size=28, bold=True, color=ACCENT)
txt(s1, "Evaluating how LLM-based routing strategies choose between local knowledge retrieval\nand real-time web search as query ambiguity increases",
    2.0, 3.65, 9.5, 0.8, size=13, color=SUBTITLE)
txt(s1, "Parin Shah  ·  Henil Patel  ·  Heta Shah", 2.0, 4.6, 9.5, 0.4, size=13, color=SUBTITLE)
txt(s1, "CS 6180  ·  Foundations of Generative AI  ·  Team 8", 2.0, 5.02, 9.5, 0.35, size=11, color=SUBTITLE)
slide_num(s1, 1)

# ── SLIDE 2: Problem & Objectives ───────────────────────────────────────────
s2 = add_slide()
accent_bar(s2)
txt(s2, "Project Introduction & Objectives", 0.7, 0.3, 10, 0.5, size=22, bold=True)
txt(s2, "Why routing is the critical decision layer in agentic RAG", 0.7, 0.82, 10, 0.35, size=13, color=SUBTITLE)

rect(s2, 0.5, 1.35, 6.1, 5.75, fill=WHITE, line=LIGHT_LINE)
txt(s2, "The Problem", 0.75, 1.5, 5.6, 0.38, size=14, bold=True, color=ACCENT)
txt(s2, "In agentic RAG, a router decides whether a query goes to a local vector database or live web search. Real-world queries are rarely clean — they are vague, context-dependent, and temporally ambiguous.",
    0.75, 1.93, 5.6, 0.9, size=11.5)
txt(s2, "Who Benefits?", 0.75, 2.95, 5.6, 0.35, size=13, bold=True)
txt(s2, "Financial analysts, compliance teams, and investors ask queries that blend static filing knowledge with live market data — one agent, two very different sources.",
    0.75, 3.35, 5.6, 0.72, size=11.5)
txt(s2, "Research Gap", 0.75, 4.18, 5.6, 0.35, size=13, bold=True)
txt(s2, "Prior RAG evaluations focus on retrieval quality, not routing decisions. No principled benchmarks exist for routing quality across ambiguity levels.",
    0.75, 4.58, 5.6, 0.65, size=11.5)
txt(s2, "Research Question", 0.75, 5.38, 5.6, 0.35, size=13, bold=True, color=ACCENT2)
txt(s2, "How does routing accuracy degrade as prompt ambiguity increases?",
    0.75, 5.78, 5.6, 0.55, size=12, bold=True)

# Right side: system diagram
rect(s2, 6.9, 1.35, 6.1, 5.75, fill=WHITE, line=LIGHT_LINE)
txt(s2, "System Architecture", 7.1, 1.5, 5.7, 0.38, size=14, bold=True, color=ACCENT)

def dbox(slide, label, sub, l, t, w, h, fc, tc):
    rect(slide, l, t, w, h, fill=fc, line=LIGHT_LINE)
    txt(slide, label, l+0.12, t+0.1, w-0.24, 0.3, size=10.5, bold=True, color=tc)
    if sub:
        txt(slide, sub, l+0.12, t+0.42, w-0.24, 0.28, size=10, color=SUBTITLE)

dbox(s2, "User Query", "", 8.4, 2.05, 3.1, 0.55, RGBColor(0xE8, 0xF0, 0xFE), ACCENT)
color_bar(s2, 9.9, 2.65, 0.05, 0.35, ACCENT)
dbox(s2, "Router Agent", "LLM-based classification", 7.9, 3.05, 4.1, 0.7, RGBColor(0xE8, 0xF0, 0xFE), ACCENT)
color_bar(s2, 8.1, 3.82, 0.05, 0.38, ACCENT2)
color_bar(s2, 11.75, 3.82, 0.05, 0.38, ACCENT3)
dbox(s2, "Local Vector DB", "ChromaDB · 22 10-Ks", 7.1, 4.23, 2.3, 0.72, RGBColor(0xE6, 0xF4, 0xEA), ACCENT2)
dbox(s2, "Web Search", "Real-time via Tavily", 10.5, 4.23, 2.3, 0.72, RGBColor(0xFE, 0xF7, 0xE0), ACCENT3)
rect(s2, 7.1, 5.15, 5.7, 1.7, fill=RGBColor(0xFF, 0xEB, 0xEB), line=RGBColor(0xF2, 0x8B, 0x82))
txt(s2, "⚠  Wrong Routing = Wrong Answer", 7.25, 5.22, 5.4, 0.35, size=11, bold=True, color=RGBColor(0xC5, 0x22, 0x1F))
txt(s2, "Routing 'stock price' to a 2025 10-K returns stale data.\nRouting 'business segments' to web misses deep filing knowledge.",
    7.25, 5.62, 5.4, 0.9, size=10.5)
slide_num(s2, 2)

# ── SLIDE 3: What We Built ──────────────────────────────────────────────────
s3 = add_slide()
accent_bar(s3, ACCENT2)
txt(s3, "Progress — What We Built", 0.7, 0.3, 10, 0.5, size=22, bold=True)
txt(s3, "6 routing strategies × 4 LLM backends × 100 labeled queries across 22 company 10-Ks", 0.7, 0.82, 12, 0.35, size=13, color=SUBTITLE)

cols = [
    ("Routing Strategies", ACCENT, [
        "Rule-Based — temporal keyword matching",
        "Always-Local — naive baseline (52% ceiling)",
        "Zero-Shot LLM — direct prompt → JSON",
        "Few-Shot LLM — 3 labeled examples in context",
        "CoT LLM — chain-of-thought before classifying",
        "LangGraph Agent — multi-node stateful workflow",
    ]),
    ("Models Evaluated", ACCENT3, [
        "Llama-3-8B-Instruct",
        "GPT-4o-mini",
        "Gemma-3-27B-IT",
        "Claude Sonnet 4.6",
        "",
        "All via OpenRouter API",
    ]),
    ("Stack & Dataset", ACCENT2, [
        "LangChain + LangGraph",
        "ChromaDB  (22 FY2025 10-Ks)",
        "Tavily real-time web search",
        "all-MiniLM-L6-v2 embeddings",
        "100 queries  ·  3 ambiguity tiers",
        "52 local  ·  48 web labels",
    ]),
]
for i, (title, color, items) in enumerate(cols):
    x = 0.5 + i * 4.28
    rect(s3, x, 1.35, 4.0, 5.75, fill=WHITE, line=LIGHT_LINE)
    color_bar(s3, x, 1.35, 4.0, 0.08, color)
    txt(s3, title, x+0.15, 1.5, 3.7, 0.38, size=13, bold=True, color=color)
    for j, item in enumerate(items):
        if item:
            txt(s3, f"• {item}", x+0.15, 2.05+j*0.82, 3.7, 0.72, size=11.5)
slide_num(s3, 3)

# ── SLIDE 4: Results — Accuracy by Tier ─────────────────────────────────────
s4 = add_slide()
accent_bar(s4, ACCENT3)
txt(s4, "Results — Routing Accuracy by Ambiguity Tier", 0.7, 0.3, 11, 0.5, size=22, bold=True)
txt(s4, "LLM-based routers decisively outperform baselines — but all degrade under ambiguity", 0.7, 0.82, 11, 0.35, size=13, color=SUBTITLE)

add_img(s4, f"{RESULTS}/accuracy_vs_tier.png", 0.4, 1.25, 7.6, 5.9)

cards = [
    (ACCENT2, "LLM Routers Win",
     "All LLM strategies reach 84% at 100% coverage. Rule-based collapses to 36% and abstains on 59% of queries."),
    (ACCENT3, "Ambiguity is the Ceiling",
     "LLM routers drop 27–37 accuracy points from Tier 1 → Tier 3. ANOVA p < 0.001 across all models."),
    (ACCENT, "Few-Shot Leads Tier 2",
     "Few-Shot hits 83% on moderately ambiguous queries — 10 pts above Zero-Shot. Best practical production gain."),
]
for i, (color, head, body) in enumerate(cards):
    yt = 1.3 + i * 1.9
    rect(s4, 8.2, yt, 4.8, 1.75, fill=WHITE, line=LIGHT_LINE)
    color_bar(s4, 8.2, yt, 0.07, 1.75, color)
    txt(s4, head, 8.4, yt+0.12, 4.4, 0.38, size=12, bold=True, color=color)
    txt(s4, body, 8.4, yt+0.55, 4.4, 1.05, size=11)
slide_num(s4, 4)

# ── SLIDE 5: Model Comparison & Failures ────────────────────────────────────
s5 = add_slide()
accent_bar(s5, ACCENT)
txt(s5, "Model Comparison & Failure Patterns", 0.7, 0.3, 11, 0.5, size=22, bold=True)
txt(s5, "Model quality drives Tier 1; ambiguity becomes the bottleneck at Tier 3", 0.7, 0.82, 11, 0.35, size=13, color=SUBTITLE)

add_img(s5, f"{RESULTS}/model_comparison.png", 0.4, 1.2, 6.4, 3.25)
add_img(s5, f"{RESULTS}/failure_patterns.png", 0.4, 4.5, 6.4, 2.75)

# Model table
rect(s5, 7.0, 1.2, 6.0, 3.25, fill=WHITE, line=LIGHT_LINE)
txt(s5, "Model Summary", 7.2, 1.32, 5.6, 0.38, size=13, bold=True)
rows = [
    ("Model",          "Overall",  "Tier 1",   "Tier 3",    True),
    ("Claude Sonnet",  "84%",      "100%",     "63–73%",    False),
    ("Gemma-3-27B",    "75–82%",   "90–98%",   "57–67%",    False),
    ("GPT-4o-mini",    "78–79%",   "98–100%",  "60–67%",    False),
    ("Llama-3-8B",     "74–77%",   "90–100%",  "57–63%",    False),
]
for i, (m, ov, t1, t3, hdr) in enumerate(rows):
    yt = 1.78 + i * 0.48
    bg = CARD_BG if hdr else WHITE
    rect(s5, 7.0, yt, 6.0, 0.46, fill=bg, line=LIGHT_LINE)
    txt(s5, m,  7.12, yt+0.08, 2.1,  0.32, size=10, bold=hdr)
    txt(s5, ov, 9.22, yt+0.08, 1.2,  0.32, size=10, bold=hdr, align=PP_ALIGN.CENTER)
    txt(s5, t1, 10.42, yt+0.08, 1.15, 0.32, size=10, bold=hdr, color=ACCENT2, align=PP_ALIGN.CENTER)
    txt(s5, t3, 11.57, yt+0.08, 1.3,  0.32, size=10, bold=hdr, color=RGBColor(0xC5, 0x22, 0x1F), align=PP_ALIGN.CENTER)

# Failure callouts
rect(s5, 7.0, 4.5, 6.0, 2.75, fill=WHITE, line=LIGHT_LINE)
txt(s5, "Top Failure Modes", 7.2, 4.62, 5.6, 0.38, size=13, bold=True)
failures = [
    ("Vague-Query Bias", "Open-ended queries default to web even when the local 10-K has the answer"),
    ("Local → Web Misroutes", "Risk factors, KPIs, strategy queries sent to web search"),
    ("Temporal Confusion", "'Early 2026' — annual filing vs. live market data ambiguity"),
]
for i, (h, b) in enumerate(failures):
    yt = 5.1 + i * 0.72
    txt(s5, f"• {h}", 7.2, yt, 5.6, 0.28, size=11, bold=True, color=ACCENT)
    txt(s5, b, 7.2, yt+0.3, 5.6, 0.36, size=10.5, color=DARK_TEXT)
slide_num(s5, 5)

# ── SLIDE 6: Next Steps ──────────────────────────────────────────────────────
s6 = add_slide()
accent_bar(s6, ACCENT2)
txt(s6, "Next Steps", 0.7, 0.3, 10, 0.5, size=22, bold=True)
txt(s6, "Two focused tasks to close the project", 0.7, 0.82, 10, 0.35, size=13, color=SUBTITLE)

tasks = [
    (ACCENT, "Linguistic Failure Analysis",
     "Examine which specific query features drive misroutes:\n\n"
     "• Missing temporal markers\n"
     "• Bare entity references (no verb, no scope)\n"
     "• Underspecified or open-ended queries\n\n"
     "Goal: characterize the residual 16–26% error rate at a linguistic level and produce interpretable findings about what makes a query hard to route.",
     "Parin Shah"),
    (ACCENT2, "Final Report",
     "Full write-up covering:\n\n"
     "• Related work & background\n"
     "• System design & methodology\n"
     "• Full experimental results\n"
     "• Limitations & future work\n\n"
     "The codebase is already reproducible. The report wraps the complete narrative around our findings.",
     "Heta Shah"),
]
for i, (color, title, body, owner) in enumerate(tasks):
    x = 0.5 + i * 6.55
    rect(s6, x, 1.35, 6.2, 5.4, fill=WHITE, line=LIGHT_LINE)
    color_bar(s6, x, 1.35, 6.2, 0.09, color)
    txt(s6, title, x+0.2, 1.52, 5.8, 0.45, size=16, bold=True, color=color)
    txt(s6, body, x+0.2, 2.1, 5.8, 3.9, size=12)
    rect(s6, x+0.2, 5.9, 5.8, 0.55, fill=CARD_BG, line=LIGHT_LINE)
    txt(s6, f"Owner: {owner}", x+0.35, 6.0, 5.5, 0.32, size=11, bold=True, color=ACCENT)

color_bar(s6, 0, 7.15, 13.33, 0.35, DARK_TEXT)
txt(s6, "Team 8  ·  Parin Shah  ·  Henil Patel  ·  Heta Shah  ·  CS 6180 · Foundations of Generative AI",
    0.5, 7.16, 12.3, 0.32, size=10, color=WHITE, align=PP_ALIGN.CENTER)
slide_num(s6, 6)

prs.save(OUT)
print("DONE:", OUT)
